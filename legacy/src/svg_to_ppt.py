#!/usr/bin/env python3
import io
from pathlib import Path

# -----------------------------
# Edit this section only
# -----------------------------
OUTPUT_PPTX = "/home/shapirma/CRELAT/output/bert_normalized_slides.pptx"
SVG_FILES = [
    "/home/shapirma/CRELAT/output/Hamlet_bert_interactions_scatter_normalized.svg",
    "/home/shapirma/CRELAT/output/Richard_II_bert_interactions_scatter_normalized.svg",
    "/home/shapirma/CRELAT/output/King_Lear_bert_interactions_scatter_normalized.svg",
    "/home/shapirma/CRELAT/output/Troilus_and_Cressida_bert_interactions_scatter_normalized.svg",
    "/home/shapirma/CRELAT/output/Twelfth_Night_bert_interactions_scatter_normalized.svg",
    "/home/shapirma/CRELAT/output/The_Merry_Wives_of_Windsor_bert_interactions_scatter_normalized.svg",
    "/home/shapirma/CRELAT/output/Much_Ado_About_Nothing_bert_interactions_scatter_normalized.svg",
    "/home/shapirma/CRELAT/output/As_You_Like_It_bert_interactions_scatter_normalized.svg",
]
# Higher DPI gives sharper images, but bigger PPT size.
RASTER_DPI = 192


def _title_from_path(path_obj):
    return path_obj.stem.replace("_", " ")


def _unique_resolved_paths(paths):
    seen = set()
    unique_paths = []
    for p in paths:
        rp = Path(p).expanduser().resolve()
        if rp not in seen:
            seen.add(rp)
            unique_paths.append(rp)
    return unique_paths


def main():
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise SystemExit(
            "Missing dependency: python-pptx. Install with `pip install python-pptx`."
        ) from exc

    try:
        import cairosvg
    except ImportError as exc:
        raise SystemExit(
            "Missing dependency: cairosvg. Install with `pip install cairosvg`."
        ) from exc

    svg_paths = _unique_resolved_paths(SVG_FILES)
    if not svg_paths:
        raise SystemExit("SVG_FILES is empty. Add at least one .svg path.")

    missing = [str(p) for p in svg_paths if not p.exists()]
    if missing:
        raise SystemExit("These files do not exist:\n- " + "\n- ".join(missing))

    invalid = [str(p) for p in svg_paths if p.suffix.lower() != ".svg"]
    if invalid:
        raise SystemExit("These files are not SVG:\n- " + "\n- ".join(invalid))

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    for svg_path in svg_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        title_shape = slide.shapes.title
        title_shape.text = _title_from_path(svg_path)

        png_bytes = cairosvg.svg2png(url=str(svg_path), dpi=RASTER_DPI)
        img_stream = io.BytesIO(png_bytes)
        pic = slide.shapes.add_picture(img_stream, 0, 0)

        slide_w = prs.slide_width
        slide_h = prs.slide_height
        margin = Inches(0.35)
        top = title_shape.top + title_shape.height + Inches(0.1)
        max_w = slide_w - (2 * margin)
        max_h = slide_h - top - margin

        scale = min(max_w / pic.width, max_h / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int(margin + (max_w - pic.width) / 2)
        pic.top = int(top + (max_h - pic.height) / 2)

    output_path = Path(OUTPUT_PPTX).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"Saved presentation: {output_path}")
    print(f"Slides created: {len(svg_paths)}")


if __name__ == "__main__":
    main()
