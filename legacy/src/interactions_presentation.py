#!/usr/bin/env python3
"""
Build a PowerPoint where each slide shows all embedding-model plots for one play (XML).

Requires: python-pptx, cairosvg
"""

from __future__ import annotations

import io
import math
from pathlib import Path

def _svg_to_png_bytes(svg_path: Path, dpi: int) -> bytes:
    try:
        import cairosvg
    except ImportError as exc:
        raise SystemExit(
            "cairosvg is required for presentations. Install with `pip install cairosvg`."
        ) from exc
    return cairosvg.svg2png(url=str(svg_path), dpi=dpi)


def build_multimodel_slides_presentation(
    slides_spec: list[dict],
    output_pptx: str | Path,
    *,
    raster_dpi: int = 192,
    slide_width_inches: float = 13.333,
    slide_height_inches: float = 7.5,
) -> Path:
    """
    Args:
        slides_spec: One entry per slide, e.g.
            [{"title": "Hamlet", "panels": [{"label": "bert-base-uncased", "svg": Path}, ...]}, ...]
        output_pptx: Destination .pptx path.
        raster_dpi: SVG rasterization resolution.

    Returns:
        Resolved path to the saved presentation.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise SystemExit(
            "python-pptx is required for presentations. Install with `pip install python-pptx`."
        ) from exc

    out = Path(output_pptx).expanduser().resolve()
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)

    margin = Inches(0.3)
    label_h = Inches(0.22)

    for spec in slides_spec:
        title_text = spec["title"]
        panels = spec["panels"]
        if not panels:
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title
        title_shape.text = title_text

        n = len(panels)
        cols = min(n, int(math.ceil(math.sqrt(n))))
        if n == 3:
            cols = 3
        rows = int(math.ceil(n / cols))

        title_bottom = title_shape.top + title_shape.height + Inches(0.08)
        avail_w = prs.slide_width - 2 * margin
        avail_h = prs.slide_height - title_bottom - margin
        cell_w = avail_w / cols
        cell_h = avail_h / rows

        for i, panel in enumerate(panels):
            svg_path = Path(panel["svg"]).expanduser().resolve()
            label = panel.get("label") or svg_path.stem
            if not svg_path.is_file():
                raise FileNotFoundError(f"Missing SVG for slide '{title_text}': {svg_path}")

            row, col = i // cols, i % cols
            cell_left = int(margin + col * cell_w)
            cell_top = int(title_bottom + row * cell_h)

            tb = slide.shapes.add_textbox(cell_left, cell_top, int(cell_w), int(label_h))
            p = tb.text_frame.paragraphs[0]
            p.text = label
            p.font.size = Pt(10)

            png_bytes = _svg_to_png_bytes(svg_path, raster_dpi)
            stream = io.BytesIO(png_bytes)
            pic = slide.shapes.add_picture(stream, cell_left, cell_top + int(label_h))

            max_w = cell_w - Inches(0.08)
            max_h = cell_h - label_h - Inches(0.08)
            scale = min(max_w / pic.width, max_h / pic.height)
            pic.width = int(pic.width * scale)
            pic.height = int(pic.height * scale)
            pic.left = cell_left + int((cell_w - pic.width) / 2)
            pic.top = cell_top + int(label_h) + int((max_h - pic.height) / 2)

    prs.save(str(out))
    return out
